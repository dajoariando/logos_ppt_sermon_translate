from pptx.util import Cm, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN

def add_page (prs, input_text, input_notes):
    """
    Adds a new blank slide to the provided Presentation object. On this new slide, 
    creates a textbox (centered vertically and horizontally within margins), places 
    the specified input text inside it, and adds speaker notes.

    Args:
        prs (Presentation): A python-pptx Presentation object.
        input_text (str): The text to display in the slide’s main textbox.
        input_notes (str): The text for the speaker notes of this slide.

    Returns:
        Presentation: The updated Presentation object with the newly added slide.
    """

    # add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # create a box in the middle of the page
    paper_width_cm = prs.slide_width
    paper_height_cm = prs.slide_height
    left = Cm(1)
    top = Cm(2)
    width = paper_width_cm - 2*left
    height = paper_height_cm - 2*top

    # add a textbox shape to the slide
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    text_frame.text = input_text

    # Optional: turn on auto-size or word wrap as needed
    # text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    # Iterate through each paragraph in the text frame
    for paragraph in text_frame.paragraphs:
        # Align the paragraph text to the left
        paragraph.alignment = PP_ALIGN.LEFT
        # Set the font size of the paragraph text to 40 points
        paragraph.font.size = Pt(40)
        # Set the font name of the paragraph text to 'Aptos (Body)'
        paragraph.font.name = 'Aptos (Body)'

    # add notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = input_notes

    return prs